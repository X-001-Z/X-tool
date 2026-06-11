from __future__ import annotations

import argparse
import io
import math
import os
import queue
import sys
import tempfile
import threading
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Iterable


VERSION = "1.5.1"
APP_TITLE = f"PDF2PPT v{VERSION}"
EMU_PER_POINT = 12700
EMU_PER_INCH = 914400
MAX_SLIDE_INCHES = 56
DEFAULT_MAX_PIXELS = 40_000_000

PRESETS = {
    "screen": (150, 85),
    "balanced": (200, 88),
    "high": (300, 92),
}

StatusCallback = Callable[[str, int | None], None]


class ConversionError(Exception):
    pass


class PasswordRequiredError(ConversionError):
    pass


class InvalidPasswordError(ConversionError):
    pass


class ConversionCancelled(ConversionError):
    pass


@dataclass
class ConversionOptions:
    dpi: int = 200
    quality: int = 88
    max_pixels: int = DEFAULT_MAX_PIXELS
    output_dir: Path | None = None
    overwrite: bool = False
    password: str | None = None

    def validate(self) -> None:
        if not 72 <= self.dpi <= 600:
            raise ConversionError("DPI 必须在 72 到 600 之间。")
        if not 50 <= self.quality <= 100:
            raise ConversionError("JPEG 质量必须在 50 到 100 之间。")
        if self.max_pixels < 1_000_000:
            raise ConversionError("像素上限不能低于 100 万。")


@dataclass
class ConversionResult:
    input_path: Path
    output_path: Path
    page_count: int
    converted_pages: int
    failed_pages: list[str] = field(default_factory=list)
    dpi_reductions: list[str] = field(default_factory=list)

    @property
    def output_size_mb(self) -> float:
        return self.output_path.stat().st_size / 1024 / 1024


def _status(callback: StatusCallback | None, text: str, percent: int | None) -> None:
    if callback:
        callback(text, percent)


def _check_cancel(cancel_event: threading.Event | None) -> None:
    if cancel_event and cancel_event.is_set():
        raise ConversionCancelled("转换已取消。")


def _unique_output_path(input_path: Path, options: ConversionOptions) -> Path:
    output_dir = options.output_dir or input_path.parent
    output_dir.mkdir(parents=True, exist_ok=True)
    candidate = output_dir / f"{input_path.stem}.pptx"
    if options.overwrite or not candidate.exists():
        return candidate

    index = 1
    while True:
        candidate = output_dir / f"{input_path.stem} ({index}).pptx"
        if not candidate.exists():
            return candidate
        index += 1


def _slide_dimensions(width_points: float, height_points: float) -> tuple[int, int]:
    if width_points <= 0 or height_points <= 0:
        raise ConversionError("PDF 页面尺寸无效。")

    max_emu = MAX_SLIDE_INCHES * EMU_PER_INCH
    width_emu = width_points * EMU_PER_POINT
    height_emu = height_points * EMU_PER_POINT
    scale = min(1.0, max_emu / max(width_emu, height_emu))
    return int(width_emu * scale), int(height_emu * scale)


def _effective_dpi(
    width_points: float, height_points: float, requested_dpi: int, max_pixels: int
) -> int:
    estimated_pixels = (
        width_points * requested_dpi / 72 * height_points * requested_dpi / 72
    )
    if estimated_pixels <= max_pixels:
        return requested_dpi
    scale = math.sqrt(max_pixels / estimated_pixels)
    return max(72, int(requested_dpi * scale))


def _picture_box(
    image_width: int,
    image_height: int,
    slide_width: int,
    slide_height: int,
) -> tuple[int, int, int, int]:
    image_ratio = image_width / image_height
    slide_ratio = slide_width / slide_height
    if image_ratio >= slide_ratio:
        width = slide_width
        height = int(width / image_ratio)
        left = 0
        top = (slide_height - height) // 2
    else:
        height = slide_height
        width = int(height * image_ratio)
        top = 0
        left = (slide_width - width) // 2
    return left, top, width, height


def inspect_pdf_password(pdf_path: str | os.PathLike[str]) -> bool:
    import pypdfium2 as pdfium

    try:
        document = pdfium.PdfDocument(str(pdf_path))
        document.close()
        return False
    except pdfium.PdfiumError as exc:
        if "password" in str(exc).lower():
            return True
        raise


def convert_pdf(
    pdf_path: str | os.PathLike[str],
    options: ConversionOptions | None = None,
    status_callback: StatusCallback | None = None,
    cancel_event: threading.Event | None = None,
) -> ConversionResult:
    import pypdfium2 as pdfium
    from pptx import Presentation

    options = options or ConversionOptions()
    options.validate()
    input_path = Path(pdf_path).expanduser().resolve()
    if not input_path.is_file():
        raise ConversionError(f"找不到文件：{input_path}")
    if input_path.suffix.lower() != ".pdf":
        raise ConversionError(f"不是 PDF 文件：{input_path.name}")

    _check_cancel(cancel_event)
    _status(status_callback, f"正在打开 {input_path.name}", 0)
    try:
        document = pdfium.PdfDocument(str(input_path), password=options.password)
    except pdfium.PdfiumError as exc:
        if "password" in str(exc).lower():
            if options.password:
                raise InvalidPasswordError("PDF 密码不正确。") from exc
            raise PasswordRequiredError("PDF 需要密码。") from exc
        raise ConversionError(f"无法打开 PDF：{exc}") from exc
    except Exception as exc:
        raise ConversionError(f"无法打开 PDF：{exc}") from exc

    try:
        page_count = len(document)
        if page_count < 1:
            raise ConversionError("PDF 没有可转换的页面。")

        first_page = document[0]
        first_width, first_height = first_page.get_size()
        first_page.close()
        slide_width, slide_height = _slide_dimensions(
            first_width, first_height
        )
        presentation = Presentation()
        presentation.slide_width = slide_width
        presentation.slide_height = slide_height
        blank_layout = presentation.slide_layouts[6]

        failed_pages: list[str] = []
        dpi_reductions: list[str] = []
        converted_pages = 0

        for page_index in range(page_count):
            _check_cancel(cancel_event)
            page_number = page_index + 1
            page = document[page_index]
            page_width, page_height = page.get_size()
            actual_dpi = _effective_dpi(
                page_width,
                page_height,
                options.dpi,
                options.max_pixels,
            )
            if actual_dpi < options.dpi:
                dpi_reductions.append(
                    f"第 {page_number} 页：{options.dpi} -> {actual_dpi} DPI"
                )

            percent = int(page_index * 90 / page_count)
            _status(
                status_callback,
                f"正在转换第 {page_number}/{page_count} 页",
                percent,
            )

            try:
                bitmap = page.render(scale=actual_dpi / 72)
                image = bitmap.to_pil().convert("RGB")
                image_buffer = io.BytesIO()
                image.save(
                    image_buffer,
                    format="JPEG",
                    quality=options.quality,
                    optimize=True,
                )
                image_bytes = image_buffer.getvalue()
                image_stream = io.BytesIO(image_bytes)
                slide = presentation.slides.add_slide(blank_layout)
                left, top, width, height = _picture_box(
                    image.width,
                    image.height,
                    slide_width,
                    slide_height,
                )
                slide.shapes.add_picture(
                    image_stream,
                    left,
                    top,
                    width=width,
                    height=height,
                )
                converted_pages += 1
            except ConversionCancelled:
                raise
            except Exception as exc:
                failed_pages.append(f"第 {page_number} 页：{exc}")
            finally:
                page.close()

        _check_cancel(cancel_event)
        if converted_pages == 0:
            details = "；".join(failed_pages[:3])
            raise ConversionError(f"没有页面转换成功。{details}")

        try:
            output_path = _unique_output_path(input_path, options)
        except OSError as exc:
            raise ConversionError(f"无法使用输出目录：{exc}") from exc
        temp_name = f".{output_path.stem}.{uuid.uuid4().hex}.tmp.pptx"
        temp_path = output_path.parent / temp_name
        _status(status_callback, "正在安全保存 PPTX", 92)
        try:
            presentation.save(str(temp_path))
            _check_cancel(cancel_event)
            os.replace(temp_path, output_path)
        except ConversionCancelled:
            temp_path.unlink(missing_ok=True)
            raise
        except Exception as exc:
            temp_path.unlink(missing_ok=True)
            raise ConversionError(f"保存 PPTX 失败：{exc}") from exc

        result = ConversionResult(
            input_path=input_path,
            output_path=output_path,
            page_count=page_count,
            converted_pages=converted_pages,
            failed_pages=failed_pages,
            dpi_reductions=dpi_reductions,
        )
        _status(
            status_callback,
            f"完成：{output_path.name}（{result.output_size_mb:.1f} MB）",
            100,
        )
        return result
    finally:
        document.close()


def _enable_console() -> None:
    if os.name != "nt":
        return
    try:
        import ctypes

        if not ctypes.windll.kernel32.AttachConsole(-1):
            ctypes.windll.kernel32.AllocConsole()
        sys.stdout = open("CONOUT$", "w", encoding="utf-8", buffering=1)
        sys.stderr = open("CONOUT$", "w", encoding="utf-8", buffering=1)
    except Exception:
        pass


def _run_cli(args: argparse.Namespace) -> int:
    _enable_console()
    output_dir = Path(args.output_dir).resolve() if args.output_dir else None
    options = ConversionOptions(
        dpi=args.dpi,
        quality=args.quality,
        max_pixels=int(args.max_megapixels * 1_000_000),
        output_dir=output_dir,
        overwrite=args.overwrite,
        password=args.password,
    )

    exit_code = 0
    for index, pdf_path in enumerate(args.files, start=1):
        print(f"[{index}/{len(args.files)}] {pdf_path}")

        def report(text: str, percent: int | None) -> None:
            suffix = f" ({percent}%)" if percent is not None else ""
            print(f"  {text}{suffix}")

        try:
            result = convert_pdf(pdf_path, options, report)
            print(f"  输出：{result.output_path}")
            if result.failed_pages:
                print(f"  警告：{len(result.failed_pages)} 页转换失败")
                exit_code = max(exit_code, 2)
        except Exception as exc:
            print(f"  错误：{exc}", file=sys.stderr)
            exit_code = 1
    return exit_code


class PDF2PPTApp:
    PRESET_LABELS = {
        "屏幕（体积小）": "screen",
        "平衡（推荐）": "balanced",
        "高清": "high",
        "自定义": "custom",
    }

    def __init__(
        self,
        root,
        initial_files: Iterable[str] = (),
        auto_start: bool = False,
    ) -> None:
        import tkinter as tk
        from tkinter import ttk

        self.tk = tk
        self.ttk = ttk
        self.root = root
        self.root.title(APP_TITLE)
        icon_path = Path(getattr(sys, "_MEIPASS", Path(__file__).resolve().parent)) / "pdf2ppt-icon.png"
        self.window_icon = None
        if icon_path.exists():
            try:
                self.window_icon = self.tk.PhotoImage(file=str(icon_path))
                self.root.iconphoto(True, self.window_icon)
            except self.tk.TclError:
                pass
        self.root.protocol("WM_DELETE_WINDOW", self._on_close)

        self.files: list[Path] = []
        self.passwords: dict[Path, str] = {}
        self.events: queue.Queue[tuple] = queue.Queue()
        self.cancel_event = threading.Event()
        self.worker: threading.Thread | None = None
        self.close_when_stopped = False
        self.last_output_dir: Path | None = None

        self.preset_var = tk.StringVar(value="平衡（推荐）")
        self.dpi_var = tk.StringVar(value="200")
        self.quality_var = tk.StringVar(value="88")
        self.output_var = tk.StringVar(value="")
        self.status_var = tk.StringVar(value="添加一个或多个 PDF 文件")
        self.summary_var = tk.StringVar(value="")

        self._build_ui()
        self._size_window_to_content()
        self._add_files(initial_files)
        self.root.after(100, self._poll_events)
        if auto_start and self.files:
            self.root.after(400, self._start)

    def _size_window_to_content(self) -> None:
        """Choose a DPI-aware size that keeps the bottom action row visible."""
        self.root.update_idletasks()
        requested_width = self.root.winfo_reqwidth()
        requested_height = self.root.winfo_reqheight()
        screen_width = self.root.winfo_screenwidth()
        screen_height = self.root.winfo_screenheight()

        width = min(max(760, requested_width + 32), max(680, screen_width - 80))
        height = min(max(620, requested_height + 32), max(540, screen_height - 100))
        min_width = min(max(680, requested_width), width)
        min_height = min(max(560, requested_height), height)

        self.root.minsize(min_width, min_height)
        left = max(0, (screen_width - width) // 2)
        top = max(0, (screen_height - height) // 2)
        self.root.geometry(f"{width}x{height}+{left}+{top}")

    def _build_ui(self) -> None:
        from tkinter import ttk

        outer = ttk.Frame(self.root, padding=16)
        outer.pack(fill="both", expand=True)

        header = ttk.Frame(outer)
        header.pack(fill="x")
        ttk.Label(header, text=APP_TITLE, font=("Segoe UI", 16, "bold")).pack(
            side="left"
        )
        ttk.Label(
            header,
            text="PDF 页面转为高清图片型 PPT",
            foreground="#555555",
        ).pack(side="left", padx=(12, 0), pady=(5, 0))

        file_frame = ttk.LabelFrame(outer, text="PDF 文件", padding=10)
        file_frame.pack(fill="both", expand=True, pady=(14, 10))
        list_area = ttk.Frame(file_frame)
        list_area.pack(fill="both", expand=True)
        self.file_list = self.tk.Listbox(
            list_area,
            selectmode="extended",
            height=9,
            font=("Segoe UI", 10),
        )
        scrollbar = ttk.Scrollbar(
            list_area, orient="vertical", command=self.file_list.yview
        )
        self.file_list.configure(yscrollcommand=scrollbar.set)
        self.file_list.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")

        file_buttons = ttk.Frame(file_frame)
        file_buttons.pack(fill="x", pady=(8, 0))
        self.add_button = ttk.Button(
            file_buttons, text="添加 PDF", command=self._browse_files
        )
        self.add_button.pack(side="left")
        self.remove_button = ttk.Button(
            file_buttons, text="移除选中", command=self._remove_selected
        )
        self.remove_button.pack(side="left", padx=6)
        self.clear_button = ttk.Button(
            file_buttons, text="清空", command=self._clear_files
        )
        self.clear_button.pack(side="left")

        settings = ttk.LabelFrame(outer, text="转换设置", padding=10)
        settings.pack(fill="x", pady=(0, 10))
        ttk.Label(settings, text="质量预设").grid(row=0, column=0, sticky="w")
        self.preset_box = ttk.Combobox(
            settings,
            textvariable=self.preset_var,
            values=list(self.PRESET_LABELS),
            state="readonly",
            width=17,
        )
        self.preset_box.grid(row=0, column=1, sticky="w", padx=(8, 20))
        self.preset_box.bind("<<ComboboxSelected>>", self._apply_preset)

        ttk.Label(settings, text="DPI").grid(row=0, column=2, sticky="w")
        self.dpi_spin = ttk.Spinbox(
            settings,
            from_=72,
            to=600,
            textvariable=self.dpi_var,
            width=7,
            command=self._mark_custom,
        )
        self.dpi_spin.grid(row=0, column=3, sticky="w", padx=(8, 20))
        self.dpi_spin.bind("<KeyRelease>", self._mark_custom)

        ttk.Label(settings, text="JPEG 质量").grid(row=0, column=4, sticky="w")
        self.quality_spin = ttk.Spinbox(
            settings,
            from_=50,
            to=100,
            textvariable=self.quality_var,
            width=7,
            command=self._mark_custom,
        )
        self.quality_spin.grid(row=0, column=5, sticky="w", padx=(8, 0))
        self.quality_spin.bind("<KeyRelease>", self._mark_custom)

        ttk.Label(settings, text="输出目录").grid(
            row=1, column=0, sticky="w", pady=(10, 0)
        )
        self.output_entry = ttk.Entry(settings, textvariable=self.output_var)
        self.output_entry.grid(
            row=1,
            column=1,
            columnspan=4,
            sticky="ew",
            padx=(8, 8),
            pady=(10, 0),
        )
        self.output_button = ttk.Button(
            settings, text="选择...", command=self._browse_output
        )
        self.output_button.grid(row=1, column=5, sticky="e", pady=(10, 0))
        settings.columnconfigure(4, weight=1)
        ttk.Label(
            settings,
            text="留空则保存到原 PDF 所在目录；已有同名文件时自动添加序号。",
            foreground="#666666",
        ).grid(row=2, column=1, columnspan=5, sticky="w", pady=(5, 0))

        self.progress = ttk.Progressbar(outer, mode="determinate", maximum=100)
        self.progress.pack(fill="x", pady=(2, 6))
        ttk.Label(outer, textvariable=self.status_var).pack(fill="x")
        ttk.Label(
            outer,
            textvariable=self.summary_var,
            foreground="#555555",
        ).pack(fill="x", pady=(2, 8))

        actions = ttk.Frame(outer)
        actions.pack(fill="x")
        self.start_button = ttk.Button(
            actions, text="开始转换", command=self._start, width=14
        )
        self.start_button.pack(side="left")
        self.cancel_button = ttk.Button(
            actions,
            text="取消",
            command=self._cancel,
            state="disabled",
            width=10,
        )
        self.cancel_button.pack(side="left", padx=8)
        self.open_button = ttk.Button(
            actions,
            text="打开输出目录",
            command=self._open_output_dir,
            state="disabled",
        )
        self.open_button.pack(side="left")
        ttk.Button(actions, text="退出", command=self._on_close, width=10).pack(
            side="right"
        )

    def _browse_files(self) -> None:
        from tkinter import filedialog

        paths = filedialog.askopenfilenames(
            title="选择 PDF 文件",
            filetypes=[("PDF 文件", "*.pdf"), ("所有文件", "*.*")],
        )
        self._add_files(paths)

    def _add_files(self, paths: Iterable[str]) -> None:
        existing = {path.resolve() for path in self.files}
        for raw_path in paths:
            path = Path(raw_path).expanduser().resolve()
            if path.is_file() and path.suffix.lower() == ".pdf" and path not in existing:
                self.files.append(path)
                existing.add(path)
                self.file_list.insert("end", str(path))
        self._refresh_file_summary()

    def _remove_selected(self) -> None:
        indexes = list(self.file_list.curselection())
        for index in reversed(indexes):
            self.file_list.delete(index)
            self.passwords.pop(self.files[index], None)
            del self.files[index]
        self._refresh_file_summary()

    def _clear_files(self) -> None:
        self.files.clear()
        self.passwords.clear()
        self.file_list.delete(0, "end")
        self._refresh_file_summary()

    def _refresh_file_summary(self) -> None:
        count = len(self.files)
        self.summary_var.set(f"已选择 {count} 个 PDF" if count else "")
        if not self.worker or not self.worker.is_alive():
            self.status_var.set("可以开始转换" if count else "添加一个或多个 PDF 文件")

    def _browse_output(self) -> None:
        from tkinter import filedialog

        selected = filedialog.askdirectory(title="选择输出目录")
        if selected:
            self.output_var.set(selected)

    def _apply_preset(self, _event=None) -> None:
        preset = self.PRESET_LABELS[self.preset_var.get()]
        if preset in PRESETS:
            dpi, quality = PRESETS[preset]
            self.dpi_var.set(str(dpi))
            self.quality_var.set(str(quality))

    def _mark_custom(self, _event=None) -> None:
        self.preset_var.set("自定义")

    def _collect_options(self) -> ConversionOptions:
        try:
            dpi = int(self.dpi_var.get().strip())
            quality = int(self.quality_var.get().strip())
        except ValueError as exc:
            raise ConversionError("DPI 和 JPEG 质量必须是整数。") from exc
        output_text = self.output_var.get().strip()
        options = ConversionOptions(
            dpi=dpi,
            quality=quality,
            output_dir=Path(output_text).expanduser().resolve() if output_text else None,
        )
        options.validate()
        return options

    def _request_passwords(self) -> bool:
        from tkinter import messagebox, simpledialog

        for path in self.files:
            try:
                needs_password = inspect_pdf_password(path)
            except Exception as exc:
                messagebox.showerror(APP_TITLE, f"无法检查 {path.name}：\n{exc}")
                return False
            if needs_password and path not in self.passwords:
                password = simpledialog.askstring(
                    APP_TITLE,
                    f"{path.name} 需要密码：",
                    show="*",
                    parent=self.root,
                )
                if password is None:
                    return False
                self.passwords[path] = password
        return True

    def _set_running(self, running: bool) -> None:
        normal = "disabled" if running else "normal"
        for widget in (
            self.add_button,
            self.remove_button,
            self.clear_button,
            self.output_button,
            self.output_entry,
            self.preset_box,
            self.dpi_spin,
            self.quality_spin,
        ):
            widget.configure(state=normal)
        self.file_list.configure(state=normal)
        if not running:
            self.preset_box.configure(state="readonly")
        self.start_button.configure(state="disabled" if running else "normal")
        self.cancel_button.configure(state="normal" if running else "disabled")

    def _start(self) -> None:
        from tkinter import messagebox

        if self.worker and self.worker.is_alive():
            return
        if not self.files:
            messagebox.showinfo(APP_TITLE, "请先添加 PDF 文件。")
            return
        try:
            options = self._collect_options()
        except ConversionError as exc:
            messagebox.showerror(APP_TITLE, str(exc))
            return
        if not self._request_passwords():
            return

        self.cancel_event.clear()
        self.progress["value"] = 0
        self.open_button.configure(state="disabled")
        self.summary_var.set("")
        self._set_running(True)
        files = list(self.files)
        passwords = dict(self.passwords)
        self.worker = threading.Thread(
            target=self._worker_main,
            args=(files, options, passwords),
            daemon=True,
        )
        self.worker.start()

    def _worker_main(
        self,
        files: list[Path],
        base_options: ConversionOptions,
        passwords: dict[Path, str],
    ) -> None:
        results: list[ConversionResult] = []
        errors: list[str] = []
        total = len(files)
        try:
            for file_index, path in enumerate(files):
                _check_cancel(self.cancel_event)
                options = ConversionOptions(
                    dpi=base_options.dpi,
                    quality=base_options.quality,
                    max_pixels=base_options.max_pixels,
                    output_dir=base_options.output_dir,
                    overwrite=base_options.overwrite,
                    password=passwords.get(path),
                )

                def report(text: str, percent: int | None) -> None:
                    if percent is None:
                        overall = None
                    else:
                        overall = int((file_index + percent / 100) * 100 / total)
                    self.events.put(("status", text, overall, file_index + 1, total))

                try:
                    results.append(
                        convert_pdf(path, options, report, self.cancel_event)
                    )
                except ConversionCancelled:
                    raise
                except Exception as exc:
                    errors.append(f"{path.name}：{exc}")
            self.events.put(("done", results, errors))
        except ConversionCancelled:
            self.events.put(("cancelled", results, errors))
        except Exception as exc:
            errors.append(str(exc))
            self.events.put(("done", results, errors))

    def _poll_events(self) -> None:
        from tkinter import messagebox

        try:
            while True:
                event = self.events.get_nowait()
                kind = event[0]
                if kind == "status":
                    _, text, percent, file_index, total = event
                    self.status_var.set(f"[{file_index}/{total}] {text}")
                    if percent is not None:
                        self.progress["value"] = percent
                elif kind == "done":
                    _, results, errors = event
                    self._finish(results, errors, cancelled=False)
                elif kind == "cancelled":
                    _, results, errors = event
                    self._finish(results, errors, cancelled=True)
        except queue.Empty:
            pass
        if self.close_when_stopped and (not self.worker or not self.worker.is_alive()):
            self.root.destroy()
            return
        self.root.after(100, self._poll_events)

    def _finish(
        self,
        results: list[ConversionResult],
        errors: list[str],
        cancelled: bool,
    ) -> None:
        from tkinter import messagebox

        self._set_running(False)
        self.worker = None
        self.progress["value"] = 100 if results and not cancelled else 0
        warnings = sum(len(result.failed_pages) for result in results)
        reductions = sum(len(result.dpi_reductions) for result in results)
        if results:
            self.last_output_dir = results[-1].output_path.parent
            self.open_button.configure(state="normal")

        if cancelled:
            self.status_var.set("转换已取消")
        elif errors:
            self.status_var.set("转换完成，但有文件失败")
        else:
            self.status_var.set("全部转换完成")

        summary = f"成功 {len(results)} 个"
        if errors:
            summary += f"，失败 {len(errors)} 个"
        if warnings:
            summary += f"，跳过 {warnings} 页"
        if reductions:
            summary += f"，{reductions} 页自动降低 DPI"
        self.summary_var.set(summary)

        if cancelled:
            messagebox.showinfo(APP_TITLE, summary + "。")
        elif errors:
            details = "\n".join(errors[:8])
            messagebox.showwarning(APP_TITLE, f"{summary}。\n\n{details}")
        else:
            messagebox.showinfo(APP_TITLE, summary + "。")

    def _cancel(self) -> None:
        self.cancel_event.set()
        self.cancel_button.configure(state="disabled")
        self.status_var.set("正在取消，请稍候...")

    def _open_output_dir(self) -> None:
        if self.last_output_dir and self.last_output_dir.exists():
            os.startfile(self.last_output_dir)

    def _on_close(self) -> None:
        from tkinter import messagebox

        if self.worker and self.worker.is_alive():
            if not messagebox.askyesno(APP_TITLE, "转换仍在进行，取消并退出吗？"):
                return
            self.cancel_event.set()
            self.close_when_stopped = True
            self.status_var.set("正在取消并退出...")
            return
        self.root.destroy()


def _run_gui(files: Iterable[str], auto_start: bool) -> int:
    import tkinter as tk

    if os.name == "nt":
        try:
            import ctypes

            ctypes.windll.shcore.SetProcessDpiAwareness(1)
        except Exception:
            pass
    root = tk.Tk()
    PDF2PPTApp(root, files, auto_start=auto_start)
    root.mainloop()
    return 0


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="将 PDF 页面转换为图片型 PPTX")
    parser.add_argument("files", nargs="*", help="一个或多个 PDF 文件")
    parser.add_argument(
        "--convert",
        action="store_true",
        help="使用命令行模式直接转换；默认打开图形界面",
    )
    parser.add_argument("--dpi", type=int, default=200, help="渲染 DPI，72-600")
    parser.add_argument(
        "--quality", type=int, default=88, help="JPEG 质量，50-100"
    )
    parser.add_argument("--output-dir", help="统一输出目录")
    parser.add_argument("--overwrite", action="store_true", help="覆盖同名 PPTX")
    parser.add_argument("--password", help="加密 PDF 的密码")
    parser.add_argument(
        "--max-megapixels",
        type=float,
        default=40,
        help="单页像素上限，默认 40",
    )
    parser.add_argument(
        "--no-auto-start",
        action="store_true",
        help="拖入 PDF 时仅载入列表，不自动开始",
    )
    parser.add_argument("--version", action="version", version=APP_TITLE)
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = _build_parser()
    args = parser.parse_args(argv)
    if args.convert:
        if not args.files:
            parser.error("命令行转换模式至少需要一个 PDF 文件")
        return _run_cli(args)
    return _run_gui(args.files, auto_start=bool(args.files) and not args.no_auto_start)


if __name__ == "__main__":
    raise SystemExit(main())
