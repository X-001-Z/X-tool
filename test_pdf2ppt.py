from __future__ import annotations

import tempfile
import threading
import unittest
from pathlib import Path
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from reportlab.pdfgen import canvas

from pdf2ppt import (
    ConversionCancelled,
    ConversionOptions,
    InvalidPasswordError,
    PasswordRequiredError,
    convert_pdf,
)


class PDF2PPTTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.folder = Path(self.temp.name)

    def tearDown(self) -> None:
        self.temp.cleanup()

    def make_pdf(self, name: str, sizes: list[tuple[int, int]]) -> Path:
        path = self.folder / name
        document = canvas.Canvas(str(path), pagesize=sizes[0])
        for index, (width, height) in enumerate(sizes):
            document.setPageSize((width, height))
            document.setFont("Helvetica", 20)
            document.drawString(30, height - 40, f"Page {index + 1}")
            document.setStrokeColorRGB(0.1, 0.4, 0.8)
            document.rect(30, 30, width - 60, height - 100, stroke=1, fill=0)
            document.showPage()
        document.save()
        return path

    def test_basic_conversion_and_no_overwrite(self) -> None:
        pdf = self.make_pdf("basic.pdf", [(595, 842)])
        options = ConversionOptions(dpi=100, quality=85)
        first = convert_pdf(pdf, options)
        second = convert_pdf(pdf, options)
        self.assertEqual(first.output_path.name, "basic.pptx")
        self.assertEqual(second.output_path.name, "basic (1).pptx")
        self.assertEqual(len(Presentation(first.output_path).slides), 1)

    def test_mixed_page_sizes_are_not_stretched(self) -> None:
        pdf = self.make_pdf("mixed.pdf", [(595, 842), (842, 595)])
        result = convert_pdf(pdf, ConversionOptions(dpi=100, quality=85))
        presentation = Presentation(result.output_path)
        self.assertEqual(len(presentation.slides), 2)
        first_picture = presentation.slides[0].shapes[0]
        second_picture = presentation.slides[1].shapes[0]
        self.assertEqual(first_picture.left, 0)
        self.assertGreater(second_picture.top, 0)
        self.assertEqual(second_picture.width, presentation.slide_width)

    def test_requested_dpi_controls_image_dimensions(self) -> None:
        pdf = self.make_pdf("dpi.pdf", [(72, 72)])
        result = convert_pdf(pdf, ConversionOptions(dpi=144, quality=85))
        with ZipFile(result.output_path) as package:
            image_name = next(
                name for name in package.namelist() if name.startswith("ppt/media/")
            )
            with package.open(image_name) as image_file:
                image = Image.open(image_file)
                self.assertEqual(image.size, (144, 144))

    def test_pixel_limit_reduces_dpi(self) -> None:
        pdf = self.make_pdf("large.pdf", [(2000, 2000)])
        result = convert_pdf(
            pdf,
            ConversionOptions(dpi=300, quality=80, max_pixels=2_000_000),
        )
        self.assertEqual(len(result.dpi_reductions), 1)

    def test_password_handling(self) -> None:
        source = self.make_pdf("plain.pdf", [(300, 300)])
        encrypted = self.folder / "encrypted.pdf"
        import pikepdf

        with pikepdf.open(source) as document:
            document.save(
                encrypted,
                encryption=pikepdf.Encryption(
                    owner="owner-secret",
                    user="user-secret",
                    R=6,
                ),
            )
        with self.assertRaises(PasswordRequiredError):
            convert_pdf(encrypted, ConversionOptions(dpi=100))
        with self.assertRaises(InvalidPasswordError):
            convert_pdf(
                encrypted,
                ConversionOptions(dpi=100, password="wrong"),
            )
        result = convert_pdf(
            encrypted,
            ConversionOptions(dpi=100, password="user-secret"),
        )
        self.assertTrue(result.output_path.exists())

    def test_pre_cancel_creates_no_output(self) -> None:
        pdf = self.make_pdf("cancel.pdf", [(595, 842)])
        event = threading.Event()
        event.set()
        with self.assertRaises(ConversionCancelled):
            convert_pdf(pdf, ConversionOptions(dpi=100), cancel_event=event)
        self.assertFalse((self.folder / "cancel.pptx").exists())


if __name__ == "__main__":
    unittest.main()
